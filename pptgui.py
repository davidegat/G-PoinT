#!/usr/bin/env python3

# G-PoinT - GPT Powered PowerPoint file and media generator.
# 2023 - gat.
# Under GPL3 License.

# Import needed libraries, se README for details.
import os
import requests
import openai
import glob
import webbrowser
import tkinter as tk
from tkinter import messagebox
from tkinter import ttk, filedialog
from tkinter import IntVar
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from gtts import gTTS

# Remember to install libraries via terminal:
# pip install os requests openai webbrowser tkinter gTTS pptx
# if you miss for some reason tkinter and glob: pip install tkinter glob

# --- MANDATORY SETTINGS ---
# Set your API key, and paths for output directory.
openai.api_key = "YOUR-API-KEY"
output_directory = "/your/desktop/or/favourite/folder/"

# --- DEFAULT SETTINGS ---
template_path = "./templates/template.pptx" # default template file path
template_folder = "./templates/" # templates folder

# Set output default language at startup.
# Language settings can be changed "on the fly" via GUI
language = "English"

# If you are NOT using English, set to False to translate prompt for DALL-E.
# This setting is changed "on the fly" via GUI.
english = True

# ---- OTHER SETTINGS ----
# Edit only if you know what to do with GPT prompts!

# Prompt for slide titles aka slide topics.
KEYPOINT_PROMPT = "Write 8 short titles of maximum 6 words on key points to be covered in a lesson about topic: {topic}. It is important that terms {topic} always appear. From now on we will interact ONLY in good {language} language."

# Prompt for slide content.
CONTENT_PROMPT = "Summarize in 6 points and using at least 10 words the most important aspects of following topic: {topic}.\nShow me ONLY the list and no other text. From now on we will interact ONLY in good {language} language."

# Picture prompt. Change last part to your favourite style, MUST be in English.
IMAGE_PROMPT = "a portrait photo of {topic}, detailed, cgi, octane, unreal"

# Define temperature value. Determines GPT randomness.
keypoint_temperature_value = 0.5
content_temperature_value = 0.7
        
# Define token length. Determines request and response lenght.
keypoint_max_tokens=1500
content_max_tokens=1500
essay_max_tokens=1500

# --- CONFIGURATION END ---

template_path = os.path.join(template_folder, "template.pptx")
image_size = "512x512"

# Get templates filenames from folder.
def get_template_filenames():
    template_files = glob.glob(os.path.join(template_folder, "*.pptx"))
    return [os.path.basename(f) for f in template_files]

# Function to update_status_bar.
def update_status_bar(text):
    status_bar.config(text=text)

# Function to generate MP3.
def generate_mp3(essay_file_path):
    mp3_file_path = os.path.splitext(essay_file_path)[0] + ".mp3"
    with open(essay_file_path, "r") as f:
        essay_text = f.read()
    language_code = language[:2].lower()
    tts = gTTS(text=essay_text, lang=language_code)
    tts.save(mp3_file_path)
    return mp3_file_path

# Function to generate PowerPoint presentation and image.
def generate_powerpoint_and_image():

    # Get topic name and number of images from input.
    topic_name = topic_entry.get()
    topic_name_en = topic_entry.get()
    num_images = int(num_images_entry.get())
    if not english:
    
        # Translate topic to English with GPT, to build DALL-E prompt later.
        prompt = f"Translate {topic_name} to English."
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=prompt,
            max_tokens=200,
            n=1,
            stop=None,
            temperature=0.5,
        )
        topic_name_en = response.choices[0].text.strip()

    # Generate key points to use as slides titles and to generate text later.
    topic_prompts = ""
    topic_name = topic_entry.get()
    prompt = KEYPOINT_PROMPT.format(topic=topic_name, language=language)
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=keypoint_max_tokens,
        n=1,
        stop=None,
        temperature=keypoint_temperature_value,
    )

    topic_prompts = response.choices[0].text.strip()
    lists = []
    for topic in topic_prompts.split("\n"):
    
        # Generate text content for each key point.
        text = ""
        prompt = CONTENT_PROMPT.format(topic=topic, language=language)
        response = openai.Completion.create(
            engine="text-davinci-003",
            prompt=prompt,
            max_tokens=content_max_tokens,
            n=1,
            stop=None,
            temperature=content_temperature_value,
        )
        text = response.choices[0].text.strip()
        lists.append((topic, text))
        
    # Create a new PowerPoint presentation from template
    # and fill it with generated key points (slides) and text.
    # See 'pptx' library documentation for further graphic customization.
    prs = Presentation(template_path)
    gpt_outputs = []  # Create an empty list for storing GPT outputs.

    for topic, text in lists:
        gpt_outputs.append((topic, text))  # Store GPT outputs.

    for topic, text in lists:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = topic
        title_shape = slide.shapes.title
        title_shape.text = topic

        # Format title font.
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.size = Pt(42)
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        for item in text.split("\n"):
            p = tf.add_paragraph()
            p.text = item.strip()
            p.level = 0
            
            # Format content font.
            p.font.bold = False
            p.font.size = Inches(0.35)
            p.font.color.rgb = RGBColor(0, 0, 0)
            
    # Save presentation.
    pptx_filename = os.path.join(output_directory, f"{topic_entry.get()}.pptx")
    prs.save(pptx_filename)     

    if essay_checkbox_var.get():
                
        # Send gpt_outputs list to GPT-3 to generate an essay.
        essay_prompt = "Please write a long detailed text in good {language} language, that considers each one of the following points:\n"
        for topic, text in gpt_outputs:
            essay_prompt += f"\n{topic}\n{text}\n"    
        response = openai.Completion.create(
            engine="text-davinci-003",
            prompt=essay_prompt,
            max_tokens=essay_max_tokens,
            n=1,
            stop=None,
            temperature=0.5,
        )  
        
    # Extract generated essay from API response
    essay = response.choices[0].text

    # Check again if essay_checkbox is checked before saving essay to a text file
    if essay_checkbox_var.get():

        # Define file path
        file_path = os.path.join(output_directory,"presentation_script.txt")

        # Write essay to a file in specified directory
        with open(file_path, "w") as f:
            f.write(essay)
    
        if generate_mp3_var.get():
            generate_mp3(file_path)

    # Generate and download an image with topic in English.
    for i in range(num_images):
        image_url = "https://api.openai.com/v1/images/generations"
        prompt = IMAGE_PROMPT.format(topic=topic_name_en)
        payload = {
            "model": "image-alpha-001",
            "prompt": prompt,
            "num_images": 1,
    	    "size": image_size,
            "response_format": "url"
        }

        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {openai.api_key}"
        }
        response = requests.post(image_url, headers=headers, json=payload)
        response.raise_for_status()
        image_url = response.json()["data"][0]["url"]
        image_content = requests.get(image_url).content

        # Save image.
        image_filename = os.path.splitext(pptx_filename)[0] + f"_image_{i+1}.png"
        with open(image_filename, "wb") as f:
            f.write(image_content)
        update_status_bar("Files generated!") 

# Function to only generate images from input.
def generate_images():
        
    # Get topic name and number of images from input.
    topic_name = topic_entry.get()
    topic_name_en = topic_entry.get()
    num_images = int(num_images_entry.get())
    if not english:
        
        # Translate topic to English with GPT, to build DALL-E prompt later.
        prompt = f"Translate {topic_name} to English."
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=prompt,
            max_tokens=200,
            n=1,
            stop=None,
            temperature=0.5,
        )
        topic_name_en = response.choices[0].text.strip()

    # Generate and download only images with topic in English.
    for i in range(num_images):
        image_url = "https://api.openai.com/v1/images/generations"
        prompt = IMAGE_PROMPT.format(topic=topic_name_en)
        payload = {
            "model": "image-alpha-001",
            "prompt": prompt,
            "num_images": 1,
 	       "size": image_size,
            "response_format": "url"
        }

        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {openai.api_key}"
        }
        response = requests.post(image_url, headers=headers, json=payload)
        response.raise_for_status()
        image_url = response.json()["data"][0]["url"]
        image_content = requests.get(image_url).content

        # Save image.
        pptx_filename = os.path.join(output_directory, f"{topic_entry.get()}.pptx")
        image_filename = os.path.splitext(pptx_filename)[0] + f"_image_{num_images + i + 1}.png"
        with open(image_filename, "wb") as f:
            f.write(image_content)
        update_status_bar(f"Generated image {num_images + i + 1} of {num_images * 2}!")
                 
# Create main window.
window = tk.Tk()
window.title("G-PoinT")

# Set window size.
window_width = 350
window_height = 500
window.geometry(f"{window_width}x{window_height}")

# Center window.
screen_width = window.winfo_screenwidth()
screen_height = window.winfo_screenheight()
x = int((screen_width/2) - (window_width/2))
y = int((screen_height/2) - (window_height/2))
window.geometry(f"+{x}+{y}")

# Label for topic entry ("Presentation topic").
topic_label = tk.Label(window, text="Presentation topic:")
topic_label.pack()

# Text entry for topic.
topic_entry = tk.Entry(window)
topic_entry.pack()

# Label for number of pictures entry ("Pictures to be generated").
num_images_label = tk.Label(window, text="How many pictures:")
num_images_label.pack()

# Text entry for number of images.
num_images_entry = tk.Entry(window)
num_images_entry.insert(0, "0")
num_images_entry.pack()

# Label for image size selection.
image_size_label = tk.Label(window, text="Picture size:")
image_size_label.pack()

# Define a variable to hold the selected image size.
image_size_var = tk.StringVar(window)
image_size_var.set("512x512")  # Set default value.

# Define available image sizes as a list of strings.
image_sizes = ["256x256", "512x512", "1024x1024"]

# Create dropdown menu for image size.
image_size_menu = tk.OptionMenu(window, image_size_var, *image_sizes)
image_size_menu.pack()

# Label for template selection.
template_label = tk.Label(window, text="Templates:")
template_label.pack()

# Define a variable to hold the selected template.
template_var = tk.StringVar(window)
template_var.set("template.pptx")  # Set default value.

# Get list of available templates.
templates = get_template_filenames()

# Create dropdown menu for templates.
template_menu = tk.OptionMenu(window, template_var, *templates)
template_menu.pack()

# Function to update template_path variable based on selected value.
def update_template_path():
    global template_path
    selected_template = template_var.get()
    template_path = os.path.join(template_folder, selected_template)

# Bind update_template_path function to dropdown menu.
template_var.trace("w", lambda *args: update_template_path())

# Function to update image_size based on selected value.
def update_image_size():
    global image_size
    image_size = image_size_var.get()

# Bind update_image_size function to dropdown menu.
image_size_var.trace("w", lambda *args: update_image_size())
separator = ttk.Separator(window, orient='horizontal')
separator.pack(fill='x', padx=10, pady=10)

generate_label = tk.Label(window, text=f"Other generation options:")
generate_label.pack(padx=2, pady=2)
essay_checkbox_var = IntVar()
essay_checkbox = tk.Checkbutton(window, text="Presentation script", variable=essay_checkbox_var)
essay_checkbox_var.set(True)
essay_checkbox.pack(padx=2, pady=2)

generate_mp3_var = tk.IntVar()

# Checkbox for MP3 generation
generate_mp3_checkbox = tk.Checkbutton(window, text=" MP3 file (script needed)", variable=generate_mp3_var)
generate_mp3_checkbox.pack()

# Button to generate PowerPoint and media.
generate_button = tk.Button(window, text="Generate PowerPoint and media", command=generate_powerpoint_and_image)
generate_button.pack(padx=10, pady=10)
separator = ttk.Separator(window, orient='horizontal')
separator.pack(fill='x', padx=10, pady=10)
generate_more_button = tk.Button(window, text="I need only pictures", command=generate_images)
generate_more_button.pack(padx=2, pady=2)

# Label to display the template path and output directory.
path_label = tk.Label(window, text=f"Templates: {template_folder}\nOutput: {output_directory}\nLanguage: {language}")
path_label.pack(padx=10, pady=10)

# Status bar label.
status_bar = tk.Label(window, text="Push and wait for notice...", bd=1, relief=tk.SUNKEN, anchor=tk.W)
status_bar.pack(side=tk.BOTTOM, fill=tk.X)

# Help message
def show_help_message():
    help_window = tk.Toplevel(window)
    help_window.title("Help")
    help_window.geometry("400x400")

    message = f"Presentation topic: fill with desired topic, word or phrase (e.g. Dolphin or Mona Lisa Story).\n\nHow many pictures: number of pictures to be generated by DALL-E. Default is 0 (no pictures).\n\nPicture size: select desired picture size from dropdown menu, choose between: 256x256, 512x512, 1024x1024.\n\nTemplate: select file from templates folder. You can add your onw template.\n\nGenerate PowerPoint and media: click to generate a full PowerPoint presentation with text, and download pics. Generated files will be saved in the output directory specified in the script.\n\nOther generation options: tick presentation essay to generate an additional file for you to read during presentation, or to include in presentation notes. Tick MP3 file, if you want the assay to be converted in MP3 format (for audio presentations or accessibility).\n\nI need only pictures: click to get additional (n) pics without PowerPoint presentation.\n\nNote: you must have a valid OpenAI API key, and set correct template and output paths in the script before running.\n\nYou can also change your default language and paths in pptgui.py file."

    text_widget = tk.Text(help_window, wrap=tk.WORD, padx=10, pady=10)
    text_widget.insert(tk.END, message)
    text_widget.pack(expand=True, fill=tk.BOTH)

    text_widget.config(state=tk.DISABLED)

    # Center help_window.
    help_window_width = 400
    help_window_height = 400
    help_window_x = window.winfo_x() + int((window_width / 2) - (help_window_width / 2))
    help_window_y = window.winfo_y() + int((window_height / 2) - (help_window_height / 2))
    help_window.geometry(f"+{help_window_x}+{help_window_y}")

# About window
def show_about_message():
    about_window = tk.Toplevel(window)
    about_window.title("About")
    about_window.geometry("400x200")
    
    message = "G-PoinT v1.5.0 by gat\n\nGenerates complete PowerPoint files and images, using OpenAI's GPT-3 API and DALL-E image generation model.\n\nhttps://github.com/davidegat/G-PoinT\n\nLicensed under the GNU General Public License v3.0"
    
    text_widget = tk.Text(about_window, wrap=tk.WORD, padx=10, pady=10)
    text_widget.insert(tk.END, message)
    text_widget.pack(expand=True, fill=tk.BOTH)

    text_widget.tag_configure("hyperlink", foreground="blue", underline=1)
    text_widget.tag_bind("hyperlink", "<Button-1>", lambda e: webbrowser.open("https://github.com/davidegat/G-PoinT"))

    start_index = text_widget.search("https://github.com/davidegat/G-PoinT", tk.END)
    end_index = text_widget.index(f"{start_index}+{len('https://github.com/davidegat/G-PoinT')}c")

    text_widget.tag_add("hyperlink", start_index, end_index)
    text_widget.config(state=tk.DISABLED)
    # Center about_window.
    about_window_width = 400
    about_window_height = 200
    about_window_x = window.winfo_x() + int((window_width / 2) - (about_window_width / 2))
    about_window_y = window.winfo_y() + int((window_height / 2) - (about_window_height / 2))
    about_window.geometry(f"+{about_window_x}+{about_window_y}")

# Language window
def open_language_window():
    # Create a new window.
    language_window = tk.Toplevel(window)
    language_window.title("Change Language")

    # Label or the language entry.
    language_label = tk.Label(language_window, text="Must match language codes: use deutsch, italian, español \nfor de, it, es, fr, ru.. or just use codes: ja, cn, en...")
    language_label.pack()

    # Text entry for language.
    language_entry = tk.Entry(language_window)
    language_entry.pack()

    # Button to update language.
    update_language_button = tk.Button(language_window, text="Update Language", command=lambda: update_language(language_entry.get(), language_window))
    update_language_button.pack()

    # Center language_window.
    language_window_width = 350
    language_window_height = 350
    language_window_x = window.winfo_x() + int((window_width / 2) - (language_window_width / 2))
    language_window_y = window.winfo_y() + int((window_height / 2) - (language_window_height / 2))
    language_window.geometry(f"+{language_window_x}+{language_window_y}")

def update_language(new_language, language_window):
    global language, english
    language = new_language
    language_code = new_language[:2].lower()  # extracts first two letters and converts to lowercase for MP3 generation.
    path_label.config(text=f"Templates: {template_folder}\nOutput: {output_directory}\nLanguage: {language}")
    if language.lower() == "english":
        english = True
    else:
        english = False
    language_window.destroy()

menu_bar = tk.Menu(window)

# Language menu.
language_menu = tk.Menu(menu_bar, tearoff=0)
menu_bar.add_command(label="Language", command=open_language_window)

# Help menu.
help_menu = tk.Menu(menu_bar, tearoff=0)
help_menu.add_command(label="Help", command=show_help_message)
help_menu.add_command(label="About", command=show_about_message)
menu_bar.add_cascade(label="Help", menu=help_menu)

# Add menu bar to main window.
window.config(menu=menu_bar)

# Start main loop.
window.mainloop()

